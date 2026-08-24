"""Automatische Prüfung erzeugter Foliensätze.

Ohne LibreOffice lässt sich kein echtes Rendering prüfen. Stattdessen laufen
hier billige, aber wirksame Strukturtests gegen die PPTX-Datei:

- ragt ein Element über den Folienrand hinaus?
- passt der Text vermutlich nicht in seinen Kasten? (Heuristik über
  durchschnittliche Zeichenbreite je Schriftgröße)
- überlappen sich zwei Textkästen?
- sind Platzhalter-Texte übrig geblieben?

Aufruf:
    python3 qa.py dist/00-orientierung.pptx        # nur Probleme
    python3 qa.py dist/00-orientierung.pptx --all  # alle Elemente auflisten
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8")

# Wörter, die im fertigen Deck nichts zu suchen haben.
PLATZHALTER_MARKER = ("PLATZHALTER", "TODO", "TBD", "XXX", "Lorem ipsum")

# Schriftgröße, die angenommen wird, wenn ein Run keine eigene setzt.
DEFAULT_FONT_PT = 18.0


def emu_to_inches(emu) -> float:
    return (emu or 0) / 914400.0


def shape_box(shape):
    return (
        emu_to_inches(shape.left),
        emu_to_inches(shape.top),
        emu_to_inches(shape.width),
        emu_to_inches(shape.height),
    )


def max_font_size(text_frame) -> float:
    """Größte tatsächlich gesetzte Schriftgröße im Rahmen."""
    sizes = [
        run.font.size.pt
        for para in text_frame.paragraphs
        for run in para.runs
        if run.font.size is not None
    ]
    return max(sizes) if sizes else DEFAULT_FONT_PT


def estimate_text_height(text: str, font_size_pt: float, width_in: float) -> float:
    """Grobe Schätzung, wie viel Höhe der Text braucht.

    Kalibriert auf Arial: die mittlere Zeichenbreite liegt bei etwa
    0.51 x Schriftgröße, bei 12pt passen also rund 11 Zeichen pro Zoll.
    Zeilenhöhe 1.3. Bewusst leicht optimistisch, damit nicht jede zweite
    Folie einen Fehlalarm auslöst – echte Überläufe werden trotzdem sichtbar.
    """
    if not text or width_in <= 0:
        return 0.0
    chars_per_inch = 11.0 * (12.0 / max(font_size_pt, 1))
    chars_per_line = max(int(width_in * chars_per_inch), 5)

    lines = 0
    for paragraph in text.split("\n"):
        if not paragraph:
            lines += 1
            continue
        lines += max((len(paragraph) + chars_per_line - 1) // chars_per_line, 1)

    return lines * (font_size_pt * 1.35) / 72.0


def overlaps(a, b, tol: float = 0.04) -> bool:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return (
        ax + aw - tol > bx
        and bx + bw - tol > ax
        and ay + ah - tol > by
        and by + bh - tol > ay
    )


def main() -> int:
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    show_all = "--all" in sys.argv

    if not args:
        print("Aufruf: python3 qa.py <datei.pptx> [--all]")
        return 2

    pptx_path = Path(args[0])
    if not pptx_path.is_absolute():
        pptx_path = Path(__file__).parent / pptx_path
    if not pptx_path.exists():
        print(f"Nicht gefunden: {pptx_path}")
        return 2

    pres = Presentation(str(pptx_path))
    sw = emu_to_inches(pres.slide_width)
    sh = emu_to_inches(pres.slide_height)

    print(f"Datei:  {pptx_path.name}")
    print(f"Format: {sw:.2f}\" x {sh:.2f}\"  ·  {len(pres.slides)} Folien")

    problem_count = 0
    placeholder_hits: list[str] = []

    for i, slide in enumerate(pres.slides, start=1):
        slide_problems: list[str] = []
        text_boxes = []  # (box, kurztext) – nur echte Textkästen

        for shape in slide.shapes:
            box = shape_box(shape)
            x, y, w, h = box
            issues = []

            if x < -0.01 or y < -0.01:
                issues.append(f"ragt oben/links hinaus ({x:.2f}, {y:.2f})")
            if x + w > sw + 0.01 or y + h > sh + 0.01:
                issues.append(
                    f"ragt unten/rechts hinaus (rechts={x + w:.2f}, unten={y + h:.2f})"
                )

            text = ""
            if shape.has_text_frame and shape.text_frame.text.strip():
                text = shape.text_frame.text
                size = max_font_size(shape.text_frame)
                needed = estimate_text_height(text, size, w)
                if needed > h + 0.08:
                    issues.append(
                        f"Text passt evtl. nicht: braucht ~{needed:.2f}\", "
                        f"Kasten {h:.2f}\" ({size:.0f}pt, {len(text)} Zeichen)"
                    )
                for marker in PLATZHALTER_MARKER:
                    if marker in text:
                        placeholder_hits.append(
                            f"Folie {i}: {text.strip().splitlines()[0][:70]}"
                        )
                        break
                text_boxes.append((box, text.replace("\n", " | ")[:50]))

            if issues:
                label = (
                    f'"{text.replace(chr(10), " | ")[:55]}"'
                    if text
                    else f"{shape.shape_type}"
                )
                slide_problems.append(
                    f"  {label}  [{w:.2f}x{h:.2f} @ {x:.2f},{y:.2f}]"
                )
                slide_problems.extend(f"      -> {it}" for it in issues)
            elif show_all:
                label = (
                    f'"{text.replace(chr(10), " | ")[:55]}"'
                    if text
                    else f"{shape.shape_type}"
                )
                slide_problems.append(
                    f"  ok  {label}  [{w:.2f}x{h:.2f} @ {x:.2f},{y:.2f}]"
                )

        # Textkästen gegeneinander prüfen
        for a_idx in range(len(text_boxes)):
            for b_idx in range(a_idx + 1, len(text_boxes)):
                (box_a, txt_a) = text_boxes[a_idx]
                (box_b, txt_b) = text_boxes[b_idx]
                if overlaps(box_a, box_b):
                    slide_problems.append(
                        f"  Überlappung: \"{txt_a}\"  <->  \"{txt_b}\""
                    )

        real_problems = [p for p in slide_problems if not p.startswith("  ok  ")]
        if real_problems:
            problem_count += len([p for p in real_problems if not p.startswith("      ->")])

        if slide_problems and (real_problems or show_all):
            print(f"\n--- Folie {i} ---")
            for line in slide_problems:
                print(line)

    print("\n" + "=" * 60)
    if placeholder_hits:
        print(f"Platzhalter noch enthalten ({len(placeholder_hits)}):")
        for hit in placeholder_hits:
            print(f"  {hit}")
    if problem_count == 0:
        print("Keine Layout-Probleme gefunden.")
    else:
        print(f"{problem_count} auffällige Stelle(n) – bitte ansehen.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

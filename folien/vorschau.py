"""Erzeugt eine HTML-Vorschau eines Foliensatzes.

Ohne installiertes Präsentationsprogramm lässt sich eine PPTX nicht rendern.
Dieses Skript zeichnet die Folien stattdessen maßstabsgetreu als HTML nach:
jede Form an ihrer Position, mit Größe, Farbe und Schriftgröße. Das reicht,
um Textüberläufe, Überlappungen und Ausreißer zu erkennen, bevor der
Foliensatz im Unterricht landet.

Aufruf:
    python3 vorschau.py dist/00-orientierung.pptx
    -> schreibt dist/00-orientierung.vorschau.html
"""

from __future__ import annotations

import html
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU_PRO_ZOLL = 914400.0
# Anzeigemaßstab: 1 Zoll entspricht so vielen Bildschirmpixeln
PX_PRO_ZOLL = 52.0


def zoll(emu) -> float:
    return (emu or 0) / EMU_PRO_ZOLL


def px(emu) -> float:
    return zoll(emu) * PX_PRO_ZOLL


def farbe_von(form) -> str | None:
    """Füllfarbe einer Form als CSS-Wert, sofern eindeutig bestimmbar."""
    try:
        fill = form.fill
        if fill.type is None:
            return None
        if fill.type == 1:  # durchgehend
            rgb = fill.fore_color.rgb
            return f"#{rgb}"
    except Exception:
        return None
    return None


def linienfarbe_von(form) -> str | None:
    try:
        line = form.line
        if line.fill.type == 1:
            return f"#{line.color.rgb}"
    except Exception:
        return None
    return None


def absatz_html(absatz) -> str:
    """Ein Absatz mit seinen Runs, Schriftgröße und Farbe."""
    teile = []
    for run in absatz.runs:
        stil = []
        f = run.font
        if f.size is not None:
            stil.append(f"font-size:{f.size.pt * PX_PRO_ZOLL / 72:.1f}px")
        if f.bold:
            stil.append("font-weight:700")
        if f.italic:
            stil.append("font-style:italic")
        try:
            if f.color is not None and f.color.type is not None:
                stil.append(f"color:#{f.color.rgb}")
        except Exception:
            pass
        if f.name:
            stil.append(f"font-family:'{f.name}',sans-serif")
        teile.append(
            f'<span style="{";".join(stil)}">{html.escape(run.text)}</span>'
        )
    inhalt = "".join(teile) or "&nbsp;"
    ausricht = {1: "center", 2: "right", 3: "justify"}.get(
        absatz.alignment.value if absatz.alignment else None, "left"
    )
    return f'<div style="text-align:{ausricht}">{inhalt}</div>'


def bild_datauri(form) -> str | None:
    """Eingebettetes Bild als data-URI, damit die Vorschau es zeigt."""
    try:
        bild = form.image
    except Exception:
        return None
    import base64
    typ = bild.content_type or "image/png"
    return f"data:{typ};base64," + base64.b64encode(bild.blob).decode("ascii")


def form_html(form) -> str:
    x, y = px(form.left), px(form.top)
    w, h = px(form.width), px(form.height)
    stil = [
        "position:absolute",
        f"left:{x:.1f}px",
        f"top:{y:.1f}px",
        f"width:{w:.1f}px",
        f"height:{h:.1f}px",
    ]

    fuell = farbe_von(form)
    linie = linienfarbe_von(form)
    if fuell:
        stil.append(f"background:{fuell}")
    if linie:
        stil.append(f"outline:1px solid {linie}")

    uri = bild_datauri(form)
    if uri:
        return (
            f'<img class="form" src="{uri}" style="{";".join(stil)};'
            f'object-fit:cover">'
        )

    inhalt = ""
    if form.has_text_frame and form.text_frame.text.strip():
        tf = form.text_frame
        vanker = {1: "center", 2: "flex-end"}.get(
            tf.vertical_anchor.value if tf.vertical_anchor else None, "flex-start"
        )
        stil += [
            "display:flex",
            "flex-direction:column",
            f"justify-content:{vanker}",
            "overflow:visible",
        ]
        inhalt = "".join(absatz_html(p) for p in tf.paragraphs)

    return f'<div class="form" style="{";".join(stil)}">{inhalt}</div>'


def main() -> int:
    if len(sys.argv) < 2:
        print("Aufruf: python3 vorschau.py <datei.pptx>")
        return 2

    pfad = Path(sys.argv[1])
    if not pfad.is_absolute():
        pfad = Path(__file__).parent / pfad
    if not pfad.exists():
        print(f"Nicht gefunden: {pfad}")
        return 2

    pres = Presentation(str(pfad))
    bw, bh = px(pres.slide_width), px(pres.slide_height)

    folien = []
    for i, folie in enumerate(pres.slides, start=1):
        formen = "".join(form_html(f) for f in folie.shapes)
        folien.append(
            f'<figure class="folie-wrap">'
            f'<figcaption>Folie {i}</figcaption>'
            f'<div class="folie" style="width:{bw:.0f}px;height:{bh:.0f}px">{formen}</div>'
            f"</figure>"
        )

    doc = f"""<!doctype html>
<html lang="de"><head><meta charset="utf-8">
<title>Vorschau: {html.escape(pfad.name)}</title>
<style>
  body {{ background:#E9EBF3; color:#272E52; font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;
          margin:0; padding:14px; }}
  .raster {{ display:flex; flex-wrap:wrap; gap:10px; }}
  h1 {{ font-size:18px; font-weight:600; margin:0 0 4px; }}
  .hinweis {{ font-size:13px; color:#7E859C; margin:0 0 24px; max-width:{bw:.0f}px; line-height:1.5; }}
  .folie-wrap {{ margin:0; }}
  figcaption {{ font-size:12px; color:#7E859C; margin-bottom:6px; font-variant-numeric:tabular-nums; }}
  .folie {{ position:relative; background:#FFFFFF; color:#434A63; overflow:hidden;
            box-shadow:0 2px 10px rgba(39,46,82,.18); font-family:Arial,Helvetica,sans-serif; }}
  .folie .form {{ line-height:1.20; }}
</style></head><body>
<h1>Vorschau: {html.escape(pfad.name)} &middot; {len(pres.slides)} Folien</h1>
<p class="hinweis">Maßstabsgetreue Nachzeichnung der Folien aus den Positionsdaten der Datei.
Schriftarten und Zeilenumbrüche können vom Präsentationsprogramm minimal abweichen –
Positionen, Größen und Textmengen stimmen.</p>
<div class="raster">{''.join(folien)}</div>
</body></html>"""

    ziel = pfad.with_suffix(".vorschau.html")
    ziel.write_text(doc, encoding="utf-8")
    print(f"{len(pres.slides)} Folien -> {ziel}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

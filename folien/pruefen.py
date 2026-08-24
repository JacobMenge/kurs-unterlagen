"""Gründliche Prüfung eines Foliensatzes.

Deckt ab, was beim Präsentieren wirklich stört:

- Text, der über seinen Kasten hinausläuft
- Elemente, die sich überlagern (auch Text über Farbfläche)
- zu kleine Schrift
- zu schwacher Kontrast zwischen Schrift- und Hintergrundfarbe
- Elemente außerhalb der Folie
- vergessene Platzhalter

Aufruf:
    python3 pruefen.py dist/00-orientierung.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

sys.stdout.reconfigure(encoding="utf-8")

PLATZHALTER = ("PLATZHALTER", "TODO", "TBD", "XXX", "Lorem ipsum")

# Mindestwerte für eine Online-Präsentation
MIN_SCHRIFT_PT = 10.0
MIN_KONTRAST = 4.5  # entspricht WCAG AA für normalen Text
MIN_KONTRAST_GROSS = 3.0  # ab 18pt bzw. 14pt fett

# Schätzung der Zeichenbreite: Arial, mittlere Breite ~0.51 x Schriftgröße
ZEICHEN_PRO_ZOLL_BEI_12PT = 11.0
ZEILENHOEHE = 1.32


def zoll(emu) -> float:
    return (emu or 0) / 914400.0


def box(shape):
    return (zoll(shape.left), zoll(shape.top), zoll(shape.width), zoll(shape.height))


def srgb(kanal: float) -> float:
    kanal /= 255.0
    return kanal / 12.92 if kanal <= 0.03928 else ((kanal + 0.055) / 1.055) ** 2.4


def luminanz(hexfarbe: str) -> float:
    h = hexfarbe.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return 0.2126 * srgb(r) + 0.7152 * srgb(g) + 0.0722 * srgb(b)


def kontrast(vorn: str, hinten: str) -> float:
    l1, l2 = luminanz(vorn), luminanz(hinten)
    hell, dunkel = max(l1, l2), min(l1, l2)
    return (hell + 0.05) / (dunkel + 0.05)


def fuellfarbe(shape) -> str | None:
    try:
        if shape.fill.type == 1:
            return str(shape.fill.fore_color.rgb)
    except Exception:
        return None
    return None


def textteile(shape):
    """Liefert (text, groesse_pt, farbe_hex, fett) je Run."""
    if not shape.has_text_frame:
        return []
    teile = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            groesse = run.font.size.pt if run.font.size is not None else 18.0
            farbe = None
            try:
                if run.font.color is not None and run.font.color.type is not None:
                    farbe = str(run.font.color.rgb)
            except Exception:
                farbe = None
            teile.append((run.text, groesse, farbe, bool(run.font.bold)))
    return teile


def hoehe_geschaetzt(text: str, pt: float, breite: float) -> float:
    if not text or breite <= 0:
        return 0.0
    pro_zoll = ZEICHEN_PRO_ZOLL_BEI_12PT * (12.0 / max(pt, 1))
    pro_zeile = max(int(breite * pro_zoll), 4)
    zeilen = 0
    for absatz in text.split("\n"):
        zeilen += max((len(absatz) + pro_zeile - 1) // pro_zeile, 1) if absatz else 1
    return zeilen * (pt * ZEILENHOEHE) / 72.0


def ueberschneidung(a, b, tol=0.03):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ox = min(ax + aw, bx + bw) - max(ax, bx)
    oy = min(ay + ah, by + bh) - max(ay, by)
    return (ox - tol, oy - tol) if (ox > tol and oy > tol) else None


def main() -> int:
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    if not args:
        print("Aufruf: python3 pruefen.py <datei.pptx>")
        return 2
    pfad = Path(args[0])
    if not pfad.is_absolute():
        pfad = Path(__file__).parent / pfad
    if not pfad.exists():
        print(f"Nicht gefunden: {pfad}")
        return 2

    pres = Presentation(str(pfad))
    bw, bh = zoll(pres.slide_width), zoll(pres.slide_height)
    print(f"Datei:  {pfad.name}")
    print(f"Format: {bw:.2f} x {bh:.2f} Zoll · {len(pres.slides)} Folien\n")

    gesamt = 0

    for nr, folie in enumerate(pres.slides, start=1):
        meldungen = []

        # Hintergrundfarbe der Folie bestimmen
        folien_bg = "FFFFFF"
        try:
            if folie.background.fill.type == 1:
                folien_bg = str(folie.background.fill.fore_color.rgb)
        except Exception:
            pass

        flaechen = []  # (box, farbe) – gefüllte Rechtecke
        texte = []  # (box, text, groesse, farbe, fett)

        for shape in folie.shapes:
            b = box(shape)
            x, y, w, h = b

            if x < -0.02 or y < -0.02 or x + w > bw + 0.02 or y + h > bh + 0.02:
                meldungen.append(
                    f"ragt aus der Folie: {shape.shape_type} bei "
                    f"({x:.2f}, {y:.2f}) Größe {w:.2f}x{h:.2f}"
                )

            farbe = fuellfarbe(shape)
            if farbe:
                flaechen.append((b, farbe))

            teile = textteile(shape)
            if teile:
                volltext = shape.text_frame.text
                max_pt = max(t[1] for t in teile)
                noetig = hoehe_geschaetzt(volltext, max_pt, w)
                if noetig > h + 0.10:
                    meldungen.append(
                        f"Text zu groß für seinen Kasten: braucht ~{noetig:.2f}\", "
                        f"hat {h:.2f}\" – „{volltext[:45].strip()}…\""
                    )
                for txt, pt, tf, fett in teile:
                    if pt < MIN_SCHRIFT_PT:
                        meldungen.append(
                            f"Schrift zu klein ({pt:.1f}pt): „{txt[:40].strip()}…\""
                        )
                texte.append((b, volltext, max_pt, teile))
                for marker in PLATZHALTER:
                    if marker in volltext:
                        meldungen.append(f"Platzhalter übrig: „{volltext[:50].strip()}…\"")

        # Kontrast: Textfarbe gegen die Fläche, auf der der Text liegt
        for tb, volltext, max_pt, teile in texte:
            hintergrund = folien_bg
            for fb, farbe in flaechen:
                u = ueberschneidung(tb, fb, tol=0.0)
                # Text liegt auf dieser Fläche, wenn er weitgehend darin liegt
                if u and u[0] > tb[2] * 0.6 and u[1] > tb[3] * 0.6:
                    hintergrund = farbe
            for txt, pt, tf, fett in teile:
                if not tf:
                    continue
                k = kontrast(tf, hintergrund)
                grenze = MIN_KONTRAST_GROSS if (pt >= 18 or (pt >= 14 and fett)) else MIN_KONTRAST
                if k < grenze:
                    meldungen.append(
                        f"Kontrast zu schwach ({k:.1f}:1, nötig {grenze}:1): "
                        f"#{tf} auf #{hintergrund}, {pt:.0f}pt – „{txt[:35].strip()}…\""
                    )

        # Text über Text
        for i in range(len(texte)):
            for j in range(i + 1, len(texte)):
                u = ueberschneidung(texte[i][0], texte[j][0])
                if u:
                    meldungen.append(
                        f"Textfelder überlagern sich ({u[0]:.2f}x{u[1]:.2f}\"): "
                        f"„{texte[i][1][:28].strip()}…\" und „{texte[j][1][:28].strip()}…\""
                    )

        if meldungen:
            gesamt += len(meldungen)
            print(f"--- Folie {nr} " + "-" * 50)
            for m in meldungen:
                print(f"  {m}")
            print()

    print("=" * 66)
    if gesamt == 0:
        print("Keine Beanstandungen.")
    else:
        print(f"{gesamt} Beanstandung(en).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
